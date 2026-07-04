from __future__ import annotations

import ctypes
import json
import os
import platform
import shutil
import stat
import subprocess
import sys
import tempfile
import time
import zipfile
from collections import deque
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any
from urllib.parse import urlsplit, urlunsplit
from uuid import uuid4

from backend.app.tools.models import ToolError

TEXT_EXTENSIONS = {".txt", ".md", ".json", ".csv", ".log"}
DOCUMENT_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".pptx"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"}
MEDIA_EXTENSIONS = IMAGE_EXTENSIONS | {".mp3", ".wav", ".m4a", ".flac", ".mp4", ".mkv", ".mov", ".avi", ".webm"}
SEARCH_EXTENSIONS = TEXT_EXTENSIONS | DOCUMENT_EXTENSIONS | MEDIA_EXTENSIONS
MAX_TEXT_BYTES = 64 * 1024
MAX_DOCUMENT_BYTES = 10 * 1024 * 1024
MAX_CLIPBOARD_CHARS = 10_000
SENSITIVE_DIRS = {
    "windows", "program files", "program files (x86)", "programdata", "recovery",
    "$recycle.bin", "system volume information", "appdata", ".ssh", ".gnupg",
    "node_modules", ".git", ".aws", ".azure", "user data", "onepassword", "bitwarden",
}
SENSITIVE_SUFFIXES = {".pem", ".key", ".pfx", ".p12", ".kdbx", ".sqlite", ".db"}
SENSITIVE_NAMES = {"credentials", "credentials.json", "id_rsa", "id_ed25519"}


@dataclass(frozen=True)
class FileReference:
    path: Path
    fingerprint: str
    expires_at: float


class FileAdapter:
    def __init__(
        self,
        root: Path,
        *,
        search_roots: list[Path] | None = None,
        reference_ttl_seconds: float = 600,
        search_timeout_seconds: float = 8,
    ) -> None:
        self.root = root.expanduser().resolve()
        self.root.mkdir(parents=True, exist_ok=True)
        self.search_roots = [path.resolve() for path in (search_roots or self._fixed_drive_roots())]
        self.reference_ttl_seconds = reference_ttl_seconds
        self.search_timeout_seconds = search_timeout_seconds
        self._references: dict[str, FileReference] = {}

    @staticmethod
    def _fixed_drive_roots() -> list[Path]:
        if sys.platform != "win32":
            return [Path.home().resolve()]
        roots: list[Path] = []
        mask = ctypes.windll.kernel32.GetLogicalDrives()
        for index in range(26):
            if mask & (1 << index):
                root = f"{chr(65 + index)}:\\"
                if ctypes.windll.kernel32.GetDriveTypeW(root) == 3:
                    roots.append(Path(root))
        return roots or [Path.home().anchor and Path(Path.home().anchor) or Path.home()]

    @staticmethod
    def _is_reparse_point(path: Path) -> bool:
        try:
            attributes = getattr(path.lstat(), "st_file_attributes", 0)
        except OSError:
            return False
        return bool(attributes & getattr(stat, "FILE_ATTRIBUTE_REPARSE_POINT", 0))

    @staticmethod
    def _is_hidden_or_system(path: Path) -> bool:
        try:
            attributes = getattr(path.stat(), "st_file_attributes", 0)
        except OSError:
            return True
        hidden = getattr(stat, "FILE_ATTRIBUTE_HIDDEN", 0)
        system = getattr(stat, "FILE_ATTRIBUTE_SYSTEM", 0)
        return bool(attributes & (hidden | system))

    def _is_sensitive_name(self, path: Path) -> bool:
        name = path.name.casefold()
        parts = path.parts
        resolved = path.resolve(strict=False)
        for root in self.search_roots:
            try:
                parts = resolved.relative_to(root).parts
                break
            except ValueError:
                continue
        return (
            name in SENSITIVE_NAMES
            or name.startswith(".env")
            or path.suffix.casefold() in SENSITIVE_SUFFIXES
            or any(part.casefold() in SENSITIVE_DIRS for part in parts)
        )

    def _allowed_root(self, path: Path) -> bool:
        resolved = path.resolve(strict=False)
        return any(resolved == root or root in resolved.parents for root in self.search_roots)

    def _validate_path(self, path: Path, *, must_exist: bool) -> Path:
        if not path.is_absolute():
            raise ToolError("TOOL_PATH_FORBIDDEN", "An absolute path on a fixed local drive is required")
        target = path.resolve(strict=False)
        if not self._allowed_root(target) or self._is_sensitive_name(target):
            raise ToolError("TOOL_PATH_FORBIDDEN", "The target is outside approved non-sensitive fixed-drive locations")
        approved_root = next(root for root in self.search_roots if target == root or root in target.parents)
        current = approved_root
        for part in target.relative_to(approved_root).parts:
            current /= part
            if current.exists() and (self._is_reparse_point(current) or self._is_hidden_or_system(current)):
                raise ToolError("TOOL_PATH_FORBIDDEN", "Hidden, system, and reparse-point paths are not allowed")
        if must_exist and not target.exists():
            raise ToolError("TOOL_FILE_NOT_FOUND", "The requested file does not exist")
        return target

    def _safe_shared_path(self, relative_path: str) -> Path:
        candidate = Path(relative_path)
        if candidate.is_absolute() or not relative_path.strip():
            raise ToolError("TOOL_PATH_FORBIDDEN", "Only relative shared-directory paths are allowed")
        target = (self.root / candidate).resolve(strict=False)
        if self.root != target and self.root not in target.parents:
            raise ToolError("TOOL_PATH_FORBIDDEN", "The requested path is outside the shared directory")
        if not target.exists():
            raise ToolError("TOOL_FILE_NOT_FOUND", "The requested file does not exist")
        return target

    @staticmethod
    def _fingerprint(path: Path) -> str:
        info = path.stat()
        return f"{info.st_dev}:{info.st_ino}:{info.st_size}:{info.st_mtime_ns}"

    @staticmethod
    def _metadata_signature(path: Path) -> tuple[int, int]:
        info = path.stat()
        return info.st_size, info.st_mtime_ns

    def _display_path(self, path: Path) -> str:
        home = Path.home().resolve()
        try:
            return f"%USERPROFILE%/{path.resolve().relative_to(home).as_posix()}"
        except ValueError:
            return str(path)

    def _search_result(self, path: Path) -> dict[str, str]:
        target = self._validate_path(path, must_exist=True)
        if not target.is_file() or target.suffix.casefold() not in SEARCH_EXTENSIONS:
            raise ToolError("TOOL_FILE_TYPE_FORBIDDEN", "The exact path is not a supported document or media file")
        file_id = str(uuid4())
        self._references[file_id] = FileReference(
            target,
            self._fingerprint(target),
            time.monotonic() + self.reference_ttl_seconds,
        )
        return {
            "file_id": file_id,
            "name": target.name,
            "path": self._display_path(target),
            "type": target.suffix.casefold().lstrip("."),
        }

    @staticmethod
    def _query_path(query: str) -> Path | None:
        value = query.strip().strip('"')
        candidate = Path(value)
        return candidate if candidate.is_absolute() else None

    def search_summary(self, query: str) -> str:
        exact = self._query_path(query)
        if exact is not None:
            target = self._validate_path(exact, must_exist=True)
            return f"Resolve exact file path: {self._display_path(target)}"
        return f"Search fixed local drives for file names containing: {query}"

    def search_details(self, query: str) -> dict[str, Any]:
        exact = self._query_path(query)
        target = self._display_path(self._validate_path(exact, must_exist=True)) if exact is not None else "Non-sensitive fixed local drives"
        return {
            "target": target,
            "operation": "search",
            "content": None,
            "content_length": None,
            "will_create_backup": False,
        }

    def search_names(self, query: str, max_results: int = 20) -> dict[str, Any]:
        exact = self._query_path(query)
        if exact is not None:
            return {"results": [self._search_result(exact)], "complete": True, "stopped_reason": None}

        normalized = query.strip().casefold()
        if not normalized or any(char in normalized for char in "*?[]\\/"):
            raise ToolError("TOOL_ARGUMENT_INVALID", "Search query must be a filename fragment or an exact absolute path")
        limit = max(1, min(max_results, 50))
        deadline = time.monotonic() + self.search_timeout_seconds
        results: list[dict[str, str]] = []
        queue = deque(self.search_roots)
        seen: set[str] = set()
        priority_names = {"desktop", "documents", "downloads", "project", "projects", "workspace", "workspaces", "桌面", "文档", "下载"}

        while queue and time.monotonic() < deadline:
            directory_path = queue.popleft()
            key = str(directory_path).casefold()
            if key in seen:
                continue
            seen.add(key)
            try:
                entries = list(os.scandir(directory_path))
            except OSError:
                continue
            for entry in entries:
                path = Path(entry.path)
                try:
                    if self._is_sensitive_name(path) or self._is_reparse_point(path) or self._is_hidden_or_system(path):
                        continue
                    if entry.is_dir(follow_symlinks=False):
                        if entry.name.casefold() in priority_names:
                            queue.appendleft(path)
                        else:
                            queue.append(path)
                        continue
                    if entry.is_file(follow_symlinks=False) and normalized in entry.name.casefold() and path.suffix.casefold() in SEARCH_EXTENSIONS:
                        results.append(self._search_result(path))
                        if entry.name.casefold() == normalized:
                            return {"results": results, "complete": False, "stopped_reason": "exact_match"}
                        if len(results) >= limit:
                            return {"results": results, "complete": False, "stopped_reason": "limit"}
                except OSError:
                    continue

        timed_out = bool(queue)
        self._purge_references()
        return {"results": results, "complete": not timed_out, "stopped_reason": "timeout" if timed_out else None}

    def _purge_references(self) -> None:
        now = time.monotonic()
        self._references = {key: value for key, value in self._references.items() if value.expires_at > now}

    def _reference(self, file_id: str) -> FileReference:
        self._purge_references()
        reference = self._references.get(file_id)
        if reference is None:
            raise ToolError("TOOL_FILE_REFERENCE_INVALID", "The file reference is missing or expired; search again")
        target = self._validate_path(reference.path, must_exist=True)
        return FileReference(target, reference.fingerprint, reference.expires_at)

    def reference_details(self, file_id: str) -> dict[str, Any]:
        reference = self._reference(file_id)
        return {"target": self._display_path(reference.path), "operation": "read", "content": None, "content_length": None, "will_create_backup": False}

    @staticmethod
    def _validate_office_archive(path: Path) -> None:
        with zipfile.ZipFile(path) as archive:
            total = 0
            for entry in archive.infolist():
                total += entry.file_size
                ratio = entry.file_size / max(1, entry.compress_size)
                if total > 50 * 1024 * 1024 or ratio > 100 or entry.filename.casefold().endswith("vbaproject.bin"):
                    raise ToolError("TOOL_FILE_ARCHIVE_UNSAFE", "The Office document archive is unsafe or contains macros")

    @staticmethod
    def _bounded(text: str) -> str:
        encoded = text.encode("utf-8")
        if len(encoded) <= MAX_TEXT_BYTES:
            return text
        return encoded[:MAX_TEXT_BYTES].decode("utf-8", errors="ignore") + "\n[truncated]"

    def read_file(self, file_id: str) -> str:
        reference = self._reference(file_id)
        path = reference.path
        if self._fingerprint(path) != reference.fingerprint:
            raise ToolError("TOOL_FILE_CHANGED", "The file changed after search; search again before reading it")
        suffix = path.suffix.casefold()
        if path.stat().st_size > MAX_DOCUMENT_BYTES:
            raise ToolError("TOOL_FILE_TOO_LARGE", "The file exceeds the 10 MiB source limit")
        try:
            if suffix in TEXT_EXTENSIONS:
                if path.stat().st_size > MAX_TEXT_BYTES:
                    raise ToolError("TOOL_FILE_TOO_LARGE", "The text file exceeds 64 KiB")
                return path.read_text(encoding="utf-8-sig")
            if suffix == ".pdf":
                from pypdf import PdfReader
                reader = PdfReader(str(path))
                if reader.is_encrypted:
                    raise ToolError("TOOL_FILE_ENCRYPTED", "Encrypted documents are not supported")
                return self._bounded("\n".join(page.extract_text() or "" for page in reader.pages))
            if suffix in {".docx", ".xlsx", ".pptx"}:
                self._validate_office_archive(path)
            if suffix == ".docx":
                from docx import Document
                return self._bounded("\n".join(paragraph.text for paragraph in Document(path).paragraphs))
            if suffix == ".xlsx":
                from openpyxl import load_workbook
                workbook = load_workbook(path, read_only=True, data_only=True)
                lines = []
                for sheet in workbook.worksheets:
                    lines.append(f"[{sheet.title}]")
                    for row in sheet.iter_rows(values_only=True):
                        lines.append("\t".join("" if value is None else str(value) for value in row))
                        if sum(len(line) for line in lines) > MAX_TEXT_BYTES:
                            break
                workbook.close()
                return self._bounded("\n".join(lines))
            if suffix == ".pptx":
                from pptx import Presentation
                lines = [shape.text for slide in Presentation(path).slides for shape in slide.shapes if hasattr(shape, "text")]
                return self._bounded("\n".join(lines))
            if suffix in MEDIA_EXTENSIONS:
                metadata: dict[str, Any] = {
                    "name": path.name,
                    "type": suffix.lstrip("."),
                    "size_bytes": path.stat().st_size,
                    "modified_at": datetime.fromtimestamp(path.stat().st_mtime, timezone.utc).isoformat(),
                }
                if suffix in IMAGE_EXTENSIONS:
                    from PIL import Image
                    with Image.open(path) as image:
                        metadata.update({"width": image.width, "height": image.height, "format": image.format})
                return json.dumps(metadata, ensure_ascii=False)
        except ToolError:
            raise
        except Exception as error:
            raise ToolError("TOOL_FILE_FORMAT_INVALID", "The file could not be safely parsed") from error
        raise ToolError("TOOL_FILE_TYPE_FORBIDDEN", "This file type is not supported")

    def read_text(self, relative_path: str) -> str:
        target = self._safe_shared_path(relative_path)
        if target.suffix.casefold() not in TEXT_EXTENSIONS:
            raise ToolError("TOOL_FILE_TYPE_FORBIDDEN", "Only approved text file types can be read")
        if target.stat().st_size > MAX_TEXT_BYTES:
            raise ToolError("TOOL_FILE_TOO_LARGE", "The file exceeds 64 KiB")
        try:
            return target.read_text(encoding="utf-8-sig")
        except UnicodeDecodeError as error:
            raise ToolError("TOOL_FILE_ENCODING_INVALID", "The file must use UTF-8 encoding") from error

    @staticmethod
    def _validate_content(content: str) -> bytes:
        encoded = content.encode("utf-8")
        if len(encoded) > MAX_TEXT_BYTES:
            raise ToolError("TOOL_CONTENT_TOO_LARGE", "Write content exceeds 64 KiB")
        return encoded

    def create_details(self, path: str, content: str) -> dict[str, Any]:
        target = self._validate_path(Path(path), must_exist=False)
        if target.suffix.casefold() not in TEXT_EXTENSIONS or not target.parent.is_dir():
            raise ToolError("TOOL_FILE_TYPE_FORBIDDEN", "Only approved text files in existing directories can be created")
        self._validate_content(content)
        return {"target": self._display_path(target), "operation": "create", "content": content, "content_length": len(content), "will_create_backup": False}

    def create_text(self, path: str, content: str) -> str:
        target = self._validate_path(Path(path), must_exist=False)
        if target.suffix.casefold() not in TEXT_EXTENSIONS or not target.parent.is_dir():
            raise ToolError("TOOL_FILE_TYPE_FORBIDDEN", "Only approved text files in existing directories can be created")
        if target.exists():
            raise ToolError("TOOL_FILE_EXISTS", "The target already exists")
        data = self._validate_content(content)
        fd, temporary = tempfile.mkstemp(prefix=".garfield-", dir=target.parent)
        try:
            with os.fdopen(fd, "wb") as handle:
                handle.write(data)
                handle.flush()
                os.fsync(handle.fileno())
            Path(temporary).rename(target)
        except Exception:
            Path(temporary).unlink(missing_ok=True)
            raise
        return f"Created {self._display_path(target)}"

    def replace_details(self, file_id: str, content: str) -> dict[str, Any]:
        reference = self._reference(file_id)
        if reference.path.suffix.casefold() not in TEXT_EXTENSIONS:
            raise ToolError("TOOL_FILE_TYPE_FORBIDDEN", "Only approved text files can be replaced")
        self._validate_content(content)
        return {"target": self._display_path(reference.path), "operation": "replace", "content": content, "content_length": len(content), "will_create_backup": True}

    def replace_text(self, file_id: str, content: str) -> str:
        reference = self._reference(file_id)
        target = reference.path
        if target.suffix.casefold() not in TEXT_EXTENSIONS:
            raise ToolError("TOOL_FILE_TYPE_FORBIDDEN", "Only approved text files can be replaced")
        if self._fingerprint(target) != reference.fingerprint:
            raise ToolError("TOOL_FILE_CHANGED", "The file changed after search; search again before replacing it")
        data = self._validate_content(content)
        stamp = datetime.now(timezone.utc).strftime("%Y%m%dT%H%M%S%fZ")
        backup = target.with_name(f"{target.name}.garfield-backup-{stamp}")
        shutil.copy2(target, backup)
        if self._fingerprint(target) != reference.fingerprint or self._metadata_signature(backup) != self._metadata_signature(target):
            backup.unlink(missing_ok=True)
            raise ToolError("TOOL_FILE_CHANGED", "The file changed while creating its backup; replacement was cancelled")
        fd, temporary = tempfile.mkstemp(prefix=".garfield-", dir=target.parent)
        try:
            with os.fdopen(fd, "wb") as handle:
                handle.write(data)
                handle.flush()
                os.fsync(handle.fileno())
            os.replace(temporary, target)
        except Exception:
            Path(temporary).unlink(missing_ok=True)
            raise
        return f"Replaced {self._display_path(target)}; backup created"


class WindowsDesktopAdapter:
    APPLICATIONS = {"notepad": ["notepad.exe"], "calculator": ["calc.exe"]}

    def __init__(self, shared_root: Path) -> None:
        self.shared_root = shared_root

    def current_time(self) -> str:
        return datetime.now().astimezone().isoformat()

    def system_info(self) -> dict[str, str]:
        return {"operating_system": platform.system(), "release": platform.release(), "architecture": platform.machine(), "python": platform.python_version()}

    def open_url(self, url: str) -> str:
        parsed = urlsplit(url.strip())
        if parsed.scheme not in {"http", "https"} or not parsed.hostname or parsed.username or parsed.password:
            raise ToolError("TOOL_URL_FORBIDDEN", "Only HTTP/HTTPS URLs without embedded credentials are allowed")
        safe_url = urlunsplit((parsed.scheme, parsed.netloc, parsed.path, parsed.query, ""))
        os.startfile(safe_url)  # type: ignore[attr-defined]
        return f"Opened {parsed.scheme}://{parsed.hostname}"

    def open_app(self, application: str) -> str:
        command = ["explorer.exe", str(self.shared_root)] if application == "explorer" else self.APPLICATIONS.get(application)
        if command is None:
            raise ToolError("TOOL_APPLICATION_FORBIDDEN", "Application is not on the allowlist")
        subprocess.Popen(command, shell=False, close_fds=True)
        return f"Opened {application}"

    def read_clipboard(self) -> str:
        if sys.platform != "win32":
            raise ToolError("TOOL_PLATFORM_UNSUPPORTED", "Clipboard tools require Windows")
        user32, kernel32 = ctypes.windll.user32, ctypes.windll.kernel32
        user32.GetClipboardData.restype = ctypes.c_void_p
        kernel32.GlobalLock.restype = ctypes.c_void_p
        if not user32.OpenClipboard(None):
            raise ToolError("TOOL_CLIPBOARD_UNAVAILABLE", "Clipboard is currently unavailable")
        try:
            handle = user32.GetClipboardData(13)
            if not handle:
                return ""
            pointer = kernel32.GlobalLock(handle)
            if not pointer:
                raise ToolError("TOOL_CLIPBOARD_UNAVAILABLE", "Clipboard text could not be read")
            try:
                text = ctypes.wstring_at(pointer)
            finally:
                kernel32.GlobalUnlock(handle)
            if len(text) > MAX_CLIPBOARD_CHARS:
                raise ToolError("TOOL_CLIPBOARD_TOO_LARGE", "Clipboard text exceeds 10,000 characters")
            return text
        finally:
            user32.CloseClipboard()

    def write_clipboard(self, text: str) -> str:
        if len(text) > MAX_CLIPBOARD_CHARS:
            raise ToolError("TOOL_CLIPBOARD_TOO_LARGE", "Clipboard text exceeds 10,000 characters")
        if sys.platform != "win32":
            raise ToolError("TOOL_PLATFORM_UNSUPPORTED", "Clipboard tools require Windows")
        subprocess.run(["powershell.exe", "-NoProfile", "-NonInteractive", "-Command", "Set-Clipboard -Value ([Console]::In.ReadToEnd())"], input=text, text=True, timeout=5, check=True, shell=False)
        return "Clipboard text updated"
